"""SID PDF Backend — FastAPI Server"""
import os
import io
import uuid
import shutil
import traceback
from pathlib import Path
from datetime import datetime, timedelta
from typing import Optional

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pypdf import PdfReader, PdfWriter
from PIL import Image

app = FastAPI(title="SID PDF API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition", "Content-Length", "Content-Type"],
)

UPLOAD_DIR = Path("/tmp/sid-pdf-uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

def cleanup_old_files():
    """Remove files older than 1 hour."""
    cutoff = datetime.now() - timedelta(hours=1)
    for f in UPLOAD_DIR.glob("*"):
        if f.is_file() and datetime.fromtimestamp(f.stat().st_mtime) < cutoff:
            f.unlink(missing_ok=True)
    for d in UPLOAD_DIR.glob("*"):
        if d.is_dir():
            try:
                if datetime.fromtimestamp(d.stat().st_mtime) < cutoff:
                    shutil.rmtree(d)
            except OSError:
                pass

@app.on_event("startup")
async def startup():
    cleanup_old_files()

@app.get("/health")
def health():
    return {"status": "ok", "service": "SID PDF API"}

# ─── Conversion ───

@app.post("/api/convert/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    """Convert PDF to Word (.docx)."""
    try:
        contents = await file.read()
        from pdf2docx import Converter

        pdf_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"
        docx_path = UPLOAD_DIR / f"{uuid.uuid4()}.docx"

        with open(pdf_path, "wb") as f:
            f.write(contents)

        cv = Converter(str(pdf_path))
        cv.convert(str(docx_path))
        cv.close()

        pdf_path.unlink(missing_ok=True)

        return FileResponse(
            docx_path,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=f"{Path(file.filename).stem}.docx",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/convert/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    """Convert PDF to Excel — extract tables."""
    try:
        contents = await file.read()
        pdf_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"

        with open(pdf_path, "wb") as f:
            f.write(contents)

        # Extract text by pages and write to simple Excel
        reader = PdfReader(str(pdf_path))
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Content"

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            ws.cell(row=i + 1, column=1, value=f"Page {i+1}")
            lines = text.split("\n")
            for j, line in enumerate(lines[:100]):
                ws.cell(row=i + 1, column=j + 2, value=line.strip())

        excel_path = UPLOAD_DIR / f"{uuid.uuid4()}.xlsx"
        wb.save(str(excel_path))
        pdf_path.unlink(missing_ok=True)

        return FileResponse(
            excel_path,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=f"{Path(file.filename).stem}.xlsx",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/convert/pdf-to-images")
async def pdf_to_images(file: UploadFile = File(...), format: str = Form("png")):
    """Convert PDF pages to images using PyMuPDF."""
    try:
        import fitz  # PyMuPDF

        contents = await file.read()
        pdf_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"

        with open(pdf_path, "wb") as f:
            f.write(contents)

        doc = fitz.open(str(pdf_path))
        output_files = []

        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=200)
            img_path = UPLOAD_DIR / f"{uuid.uuid4()}.{format}"
            pix.save(str(img_path))
            output_files.append(img_path)

        doc.close()
        pdf_path.unlink(missing_ok=True)

        if len(output_files) == 1:
            return FileResponse(
                output_files[0],
                media_type=f"image/{format}",
                filename=f"{Path(file.filename).stem}.{format}",
            )

        # ZIP multiple images
        import zipfile

        zip_path = UPLOAD_DIR / f"{uuid.uuid4()}.zip"
        with zipfile.ZipFile(str(zip_path), "w") as zf:
            for i, img in enumerate(output_files):
                zf.write(str(img), f"page_{i+1}.{format}")
                img.unlink(missing_ok=True)

        return FileResponse(
            zip_path,
            media_type="application/zip",
            filename=f"{Path(file.filename).stem}_images.zip",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/convert/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    """Convert PDF to PowerPoint — extract pages as slides with images."""
    try:
        import fitz
        from pptx import Presentation
        from pptx.util import Inches

        contents = await file.read()
        pdf_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"

        with open(pdf_path, "wb") as f:
            f.write(contents)

        doc = fitz.open(str(pdf_path))
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # Blank layout

        for page_num in range(len(doc)):
            page = doc[page_num]
            pix = page.get_pixmap(dpi=150)
            img_path = UPLOAD_DIR / f"{uuid.uuid4()}.png"
            pix.save(str(img_path))

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(img_path),
                Inches(0), Inches(0),
                prs.slide_width, prs.slide_height,
            )
            img_path.unlink(missing_ok=True)

        pptx_path = UPLOAD_DIR / f"{uuid.uuid4()}.pptx"
        prs.save(str(pptx_path))
        doc.close()
        pdf_path.unlink(missing_ok=True)

        return FileResponse(
            pptx_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"{Path(file.filename).stem}.pptx",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/convert/office-to-pdf")
async def office_to_pdf(file: UploadFile = File(...)):
    """Convert Office files to PDF. Supports DOCX and images."""
    try:
        contents = await file.read()
        ext = Path(file.filename).suffix.lower()

        input_path = UPLOAD_DIR / f"{uuid.uuid4()}{ext}"
        pdf_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"

        with open(input_path, "wb") as f:
            f.write(contents)

        if ext == ".docx":
            # Use python-docx + reportlab for reliable conversion
            from docx import Document as DocxDocument
            from docx.oxml.ns import qn
            from reportlab.pdfgen import canvas as rl_canvas
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.units import mm

            doc = DocxDocument(str(input_path))
            c = rl_canvas.Canvas(str(pdf_path), pagesize=A4)
            width, height = A4
            y = height - 40
            left_margin = 40
            right_margin = width - 40
            line_height = 14
            font_name = "Helvetica"
            font_size = 9

            def draw_text(text, bold=False, size=None, color=None):
                nonlocal y
                sz = size or font_size
                fn = "Helvetica-Bold" if bold else font_name
                c.setFont(fn, sz)
                if color:
                    c.setFillColorRGB(*color)
                else:
                    c.setFillColorRGB(0, 0, 0)

                words = text.split()
                line = ""
                for word in words:
                    test = f"{line} {word}".strip()
                    if c.stringWidth(test, fn, sz) < right_margin - left_margin:
                        line = test
                    else:
                        c.drawString(left_margin, y, line)
                        y -= line_height
                        line = word
                        if y < 40:
                            c.showPage()
                            c.setFont(fn, sz)
                            y = height - 40
                if line:
                    c.drawString(left_margin, y, line)
                    y -= line_height
                if y < 40:
                    c.showPage()
                    y = height - 40

            def draw_line(y_pos):
                c.setStrokeColorRGB(0.7, 0.7, 0.7)
                c.line(left_margin, y_pos, right_margin, y_pos)

            # Iterate through body elements in order (paragraphs + tables)
            body = doc.element.body
            for child in body:
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                
                if tag == 'p':
                    # Paragraph
                    para_text = ""
                    for run in child.findall('.//' + qn('w:r')):
                        t = run.find(qn('w:t'))
                        if t is not None and t.text:
                            para_text += t.text
                    
                    text = para_text.strip()
                    if not text:
                        y -= 6
                        if y < 40:
                            c.showPage()
                            y = height - 40
                        continue
                    
                    # Check if bold
                    is_bold = False
                    for run in child.findall('.//' + qn('w:r')):
                        rpr = run.find(qn('w:rPr'))
                        if rpr is not None:
                            b = rpr.find(qn('w:b'))
                            if b is not None:
                                is_bold = True
                                break
                    
                    sz = 11 if is_bold else font_size
                    draw_text(text, bold=is_bold, size=sz)
                    y -= 2  # spacing after paragraph
                
                elif tag == 'tbl':
                    # Table
                    y -= 6
                    if y < 60:
                        c.showPage()
                        y = height - 40
                    
                    rows = child.findall('.//' + qn('w:tr'))
                    if not rows:
                        continue
                    
                    # Calculate column widths
                    num_cols = 0
                    for row in rows:
                        cells = row.findall(qn('w:tc'))
                        num_cols = max(num_cols, len(cells))
                    
                    if num_cols == 0:
                        continue
                    
                    col_width = (right_margin - left_margin) / num_cols
                    
                    for row_idx, row in enumerate(rows):
                        cells = row.findall(qn('w:tc'))
                        cell_texts = []
                        cell_bolds = []
                        
                        for cell in cells:
                            cell_text = ""
                            cell_bold = False
                            for p in cell.findall('.//' + qn('w:p')):
                                for run in p.findall('.//' + qn('w:r')):
                                    t = run.find(qn('w:t'))
                                    if t is not None and t.text:
                                        cell_text += t.text
                                    rpr = run.find(qn('w:rPr'))
                                    if rpr is not None and rpr.find(qn('w:b')) is not None:
                                        cell_bold = True
                            cell_texts.append(cell_text.strip())
                            cell_bolds.append(cell_bold)
                        
                        # Draw row background for header
                        row_top = y
                        if row_idx == 0:
                            c.setFillColorRGB(0.9, 0.9, 0.95)
                            c.rect(left_margin, y - 12, right_margin - left_margin, 16, fill=1, stroke=0)
                            c.setFillColorRGB(0, 0, 0)
                        
                        # Draw cell text
                        c.setFont("Helvetica-Bold" if row_idx == 0 else font_name, 8)
                        for col, (ct, cb) in enumerate(zip(cell_texts, cell_bolds)):
                            x = left_margin + col * col_width + 2
                            cell_w = col_width - 4
                            # Truncate text to fit cell
                            display_text = ct
                            while display_text and c.stringWidth(display_text, c._fontname, c._fontsize) > cell_w:
                                display_text = display_text[:-1]
                            c.drawString(x, y - 10, display_text)
                        
                        y -= 16
                        # Draw row separator
                        draw_line(y + 2)
                        
                        if y < 40:
                            c.showPage()
                            y = height - 40
                    
                    y -= 6  # spacing after table

            c.save()
        elif ext in (".jpg", ".jpeg", ".png"):
            img = Image.open(str(input_path)).convert("RGB")
            img.save(str(pdf_path), "PDF")
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported format: {ext}")

        input_path.unlink(missing_ok=True)

        return FileResponse(
            pdf_path,
            media_type="application/pdf",
            filename=f"{Path(file.filename).stem}.pdf",
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"{type(e).__name__}: {str(e)}")


# ─── Compression ───

@app.post("/api/compress")
async def compress_pdf(file: UploadFile = File(...), quality: str = Form("medium")):
    """Compress PDF using pypdf compression."""
    try:
        quality_map = {"low": 50, "medium": 30, "high": 10}
        factor = quality_map.get(quality, 30)

        contents = await file.read()
        reader = PdfReader(io.BytesIO(contents))
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        # Compress after adding pages
        for page in writer.pages:
            page.compress_content_streams()

        output = io.BytesIO()
        writer.write(output)
        output.seek(0)

        compressed_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"
        with open(compressed_path, "wb") as f:
            f.write(output.read())

        return FileResponse(
            compressed_path,
            media_type="application/pdf",
            filename=f"{Path(file.filename).stem}_compressed.pdf",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── Merge PDFs ───

@app.post("/api/merge")
async def merge_pdfs(files: list[UploadFile] = File(...)):
    """Merge multiple PDFs into one."""
    try:
        writer = PdfWriter()
        temp_files = []

        for f in files:
            contents = await f.read()
            temp_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"
            with open(temp_path, "wb") as tf:
                tf.write(contents)
            reader = PdfReader(str(temp_path))
            writer.append_pages_from_reader(reader)
            temp_files.append(temp_path)

        output_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"
        writer.write(str(output_path))

        for tf in temp_files:
            tf.unlink(missing_ok=True)

        return FileResponse(
            output_path,
            media_type="application/pdf",
            filename="merged.pdf",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── Split PDF ───

@app.post("/api/split")
async def split_pdf(file: UploadFile = File(...), pages: str = Form("")):
    """Split PDF by page ranges (e.g. '1-3,5,7-9')."""
    try:
        contents = await file.read()
        reader = PdfReader(io.BytesIO(contents))
        total = len(reader.pages)

        # Parse page ranges
        page_nums = set()
        if not pages.strip():
            # Split every page
            page_nums = set(range(total))
        else:
            for part in pages.split(","):
                part = part.strip()
                if "-" in part:
                    start, end = part.split("-", 1)
                    start = max(0, int(start.strip()) - 1)
                    end = min(total, int(end.strip()))
                    page_nums.update(range(start, end))
                else:
                    p = int(part.strip()) - 1
                    if 0 <= p < total:
                        page_nums.add(p)

        if not page_nums:
            raise HTTPException(status_code=400, detail="No valid pages specified")

        output_files = []
        for p in sorted(page_nums):
            writer = PdfWriter()
            writer.add_page(reader.pages[p])
            out_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"
            writer.write(str(out_path))
            output_files.append((p + 1, out_path))

        if len(output_files) == 1:
            return FileResponse(
                output_files[0][1],
                media_type="application/pdf",
                filename=f"{Path(file.filename).stem}_page_{output_files[0][0]}.pdf",
            )

        # ZIP multiple files
        import zipfile

        zip_path = UPLOAD_DIR / f"{uuid.uuid4()}.zip"
        with zipfile.ZipFile(str(zip_path), "w") as zf:
            for page_num, fp in output_files:
                zf.write(str(fp), f"page_{page_num}.pdf")
                fp.unlink(missing_ok=True)

        return FileResponse(
            zip_path,
            media_type="application/zip",
            filename=f"{Path(file.filename).stem}_split.zip",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── Protect / Unlock ───

@app.post("/api/protect")
async def protect_pdf(
    file: UploadFile = File(...),
    password: str = Form(...),
):
    """Encrypt PDF with password."""
    try:
        contents = await file.read()
        reader = PdfReader(io.BytesIO(contents))
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        writer.encrypt(password)

        output_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"
        writer.write(str(output_path))

        return FileResponse(
            output_path,
            media_type="application/pdf",
            filename=f"{Path(file.filename).stem}_protected.pdf",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/unlock")
async def unlock_pdf(
    file: UploadFile = File(...),
    password: str = Form(""),
):
    """Remove password from PDF."""
    try:
        contents = await file.read()
        reader = PdfReader(io.BytesIO(contents))

        if reader.is_encrypted:
            try:
                reader.decrypt(password)
            except Exception:
                raise HTTPException(status_code=400, detail="Password salah atau file tidak bisa didekripsi")

        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        output_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"
        writer.write(str(output_path))

        return FileResponse(
            output_path,
            media_type="application/pdf",
            filename=f"{Path(file.filename).stem}_unlocked.pdf",
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── Watermark ───

@app.post("/api/watermark")
async def watermark_pdf(
    file: UploadFile = File(...),
    text: str = Form(""),
    opacity: float = Form(0.3),
    font_size: int = Form(60),
    rotation: float = Form(0),
):
    """Add text watermark to PDF using PyMuPDF."""
    try:
        import fitz

        # Clamp rotation to valid values (PyMuPDF only allows 0, 90, 180, 270)
        valid_rotations = [0, 90, 180, 270]
        clamped_rotation = min(valid_rotations, key=lambda x: abs(x - int(rotation)))

        contents = await file.read()
        pdf_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"

        with open(pdf_path, "wb") as f:
            f.write(contents)

        doc = fitz.open(str(pdf_path))

        for page in doc:
            rect = page.rect
            # Center of page
            x0 = rect.width / 2
            y0 = rect.height / 2

            # Create text with rotation
            page.insert_text(
                fitz.Point(x0, y0),
                text,
                fontsize=font_size,
                color=(0.5, 0.5, 0.5),  # gray
                rotate=clamped_rotation,
                overlay=True,
            )

        output_path = UPLOAD_DIR / f"{uuid.uuid4()}.pdf"
        doc.save(str(output_path))
        doc.close()
        pdf_path.unlink(missing_ok=True)

        return FileResponse(
            output_path,
            media_type="application/pdf",
            filename=f"{Path(file.filename).stem}_watermarked.pdf",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
